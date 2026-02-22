import os
from pptx import Presentation

# 定義路徑
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DOWNLOADS_DIR = os.path.join(BASE_DIR, "downloads")
# 新增：模板資料夾與檔案路徑
TEMPLATES_DIR = os.path.join(BASE_DIR, "templates")
TEMPLATE_FILE = os.path.join(TEMPLATES_DIR, "defense_template.pptx")

# 確保下載目錄存在
os.makedirs(DOWNLOADS_DIR, exist_ok=True)

def replace_text_in_slide(slide, replacements):
    """
    遍歷投影片中的所有形狀，尋找並替換文字。
    解決 PPT 底層會把同一個單字切碎 (Runs) 導致無法匹配的問題。
    """
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
            
        for paragraph in shape.text_frame.paragraphs:
            # 1. 把整個段落所有切碎的片段，先無縫接軌拼成一句完整的話
            full_text = "".join([run.text for run in paragraph.runs])
            
            # 2. 檢查這整句話裡有沒有需要替換的變數
            needs_replace = False
            for placeholder in replacements.keys():
                if placeholder in full_text:
                    needs_replace = True
                    break
            
            # 3. 如果有找到變數，就整段進行替換
            if needs_replace:
                for placeholder, value in replacements.items():
                    full_text = full_text.replace(placeholder, str(value or ""))
                    
                # 4. 為了保留你的字體排版格式，我們把替換完的完整句子塞回第一個片段，並把多餘的片段清空
                if len(paragraph.runs) > 0:
                    paragraph.runs[0].text = full_text
                    for i in range(1, len(paragraph.runs)):
                        paragraph.runs[i].text = ""

def generate_ppt(payload, log_id: int) -> str:
    """
    讀取模板 PPTX，替換其中的佔位符資料，並產出新檔案。
    """
    # 1. 檢查模板檔案是否存在
    if not os.path.exists(TEMPLATE_FILE):
        raise FileNotFoundError(f"找不到模板檔案，請確認路徑：{TEMPLATE_FILE}")

    # 2. 載入模板簡報
    prs = Presentation(TEMPLATE_FILE)
    
    # 3. 取得第一張投影片 (通常模板只有一張)
    slide = prs.slides[0]
    
    # 4. 準備要替換的資料對照表 (佔位符 -> 真實資料)
    # 先處理口試委員清單，加上縮排符號
    committee_text = "\n".join([f"    {c}" for c in payload.committee_members])
    
    replacements = {
        "{{student_name}}": payload.student_name,
        "{{student_id}}": payload.student_id,
        "{{thesis_title_zh}}": payload.thesis_title_zh,
        "{{thesis_title_en}}": payload.thesis_title_en,
        "{{advisor_full_text}}": payload.advisor_full_text,
        "{{defense_date_text}}": payload.defense_date_text,
        "{{defense_time_text}}": payload.defense_time_text,
        "{{location_full_text}}": payload.location_full_text,
        "{{committee_members_list}}": committee_text
    }
    
    # 5. 執行替換邏輯
    replace_text_in_slide(slide, replacements)
    
    # 6. 存檔
    filename = f"defense_{payload.student_id}_{log_id}.pptx"
    file_path = os.path.join(DOWNLOADS_DIR, filename)
    prs.save(file_path)
    
    print(f"✅ PPT 生成成功：{file_path}")
    # 回傳生成的檔案名稱
    return filename